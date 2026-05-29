"""Generate architecture-and-stack.pptx — 3 slide deck."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

NAVY = RGBColor(0x0F, 0x1E, 0x3D)
ACCENT = RGBColor(0x2E, 0x86, 0xDE)
LIGHT = RGBColor(0xF4, 0xF6, 0xFA)
DARK = RGBColor(0x1A, 0x1A, 0x1A)
GREY = RGBColor(0x55, 0x5B, 0x6E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height

blank = prs.slide_layouts[6]


def add_rect(slide, x, y, w, h, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
    shp.shadow.inherit = False
    return shp


def add_text(slide, x, y, w, h, text, size=18, bold=False, color=DARK, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = "Calibri"
    return tb


def slide_header(slide, title, subtitle=None):
    add_rect(slide, 0, 0, SW, Inches(0.9), NAVY)
    add_rect(slide, 0, Inches(0.9), SW, Inches(0.08), ACCENT)
    add_text(slide, Inches(0.5), Inches(0.18), Inches(12), Inches(0.6),
             title, size=28, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, Inches(0.5), Inches(1.05), Inches(12), Inches(0.4),
                 subtitle, size=14, color=GREY)


# ---------------- Slide 1: Architecture ----------------
s1 = prs.slides.add_slide(blank)
slide_header(s1, "Invoicing Platform — Architecture",
             "Server-rendered Next.js monolith · Postgres · On-demand xlsx")

# Layer stack (5 boxes vertical)
layers = [
    ("Browser (Admin / SDM)", "React 19 + Tailwind + shadcn/ui", ACCENT),
    ("Next.js 15 — App Router", "RSC + Server Actions · NextAuth v5 (Entra ID)", NAVY),
    ("Domain Services", "Rate resolution · Assignment guards · Invoice engine (ExcelJS)", NAVY),
    ("Prisma Client (singleton)", "Type-safe queries · Migration history checked-in", NAVY),
    ("PostgreSQL 15+   |   Object Storage (Blob / S3)", "Partial-unique indexes · Blob URLs only", GREY),
]
top = Inches(1.7)
box_h = Inches(0.85)
gap = Inches(0.12)
box_w = Inches(7.5)
left = Inches(0.5)
for title, sub, col in layers:
    add_rect(s1, left, top, box_w, box_h, col)
    add_text(s1, left + Inches(0.2), top + Inches(0.08), box_w - Inches(0.4), Inches(0.4),
             title, size=15, bold=True, color=WHITE)
    add_text(s1, left + Inches(0.2), top + Inches(0.45), box_w - Inches(0.4), Inches(0.4),
             sub, size=11, color=LIGHT)
    top += box_h + gap

# Right column — invariants
rx = Inches(8.3)
add_rect(s1, rx, Inches(1.7), Inches(4.5), Inches(5.2), LIGHT)
add_text(s1, rx + Inches(0.2), Inches(1.8), Inches(4.2), Inches(0.4),
         "Key Invariants", size=16, bold=True, color=NAVY)
inv_lines = [
    "• DEDICATED single-account: app + DB partial unique index",
    "• SDM scoping enforced on every query (UserAccountAccess)",
    "• Immutable updates — no in-place mutation",
    "• Rate cells nullable — fill commercials over time",
    "• DATE-only timesheets at launch",
    "",
    "Invoice Flow",
    "(account, period) → Org → template → Assignments × Entries",
    "→ AccountRate join (band + sub-cat + SLA) → ExcelJS → Blob",
]
add_text(s1, rx + Inches(0.2), Inches(2.3), Inches(4.2), Inches(4.5),
         "\n".join(inv_lines), size=11, color=DARK)

# ---------------- Slide 2: Tech Stack ----------------
s2 = prs.slides.add_slide(blank)
slide_header(s2, "Tech Stack",
             "Pinned versions · pnpm-managed · TypeScript end-to-end")

# 4x3 grid of pills
stack = [
    ("Next.js 15", "App Router + RSC"),
    ("React 19", "Server-first UI"),
    ("TypeScript 5.7", "End-to-end types"),
    ("Tailwind 3.4", "Utility-first CSS"),
    ("PostgreSQL 15+", "Decimals + partial idx"),
    ("Prisma 6.2", "Type-safe ORM"),
    ("NextAuth v5", "Entra ID OIDC"),
    ("Zod 3.24", "Boundary validation"),
    ("ExcelJS", "Template xlsx render"),
    ("Vitest 2.1", "Unit tests"),
    ("Playwright 1.50", "E2E flows"),
    ("pnpm 10.33", "Deterministic installs"),
]
cols, rows = 4, 3
grid_left = Inches(0.5)
grid_top = Inches(1.7)
cell_w = Inches(3.05)
cell_h = Inches(1.2)
cell_gap_x = Inches(0.12)
cell_gap_y = Inches(0.18)

for i, (name, desc) in enumerate(stack):
    r, c = divmod(i, cols)
    x = grid_left + c * (cell_w + cell_gap_x)
    y = grid_top + r * (cell_h + cell_gap_y)
    add_rect(s2, x, y, cell_w, cell_h, LIGHT)
    add_rect(s2, x, y, Inches(0.12), cell_h, ACCENT)
    add_text(s2, x + Inches(0.25), y + Inches(0.18), cell_w - Inches(0.3), Inches(0.4),
             name, size=14, bold=True, color=NAVY)
    add_text(s2, x + Inches(0.25), y + Inches(0.6), cell_w - Inches(0.3), Inches(0.5),
             desc, size=11, color=GREY)

# Footer band
add_rect(s2, Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.7), NAVY)
add_text(s2, Inches(0.7), Inches(6.5), Inches(12), Inches(0.5),
         "Stateless workers · DB-backed sessions · Object storage for blobs · Validate at every boundary",
         size=13, bold=True, color=WHITE)

# ---------------- Slide 3: Scalability ----------------
s3 = prs.slides.add_slide(blank)
slide_header(s3, "Scalability",
             "Compute · Database · Storage · Schema")

quads = [
    ("Compute", [
        "Stateless Next.js workers — horizontal scale",
        "No sticky session (DB-backed)",
        "Edge cache static assets",
        "Invoice gen → job queue when p99 > 5s",
    ]),
    ("Database", [
        "Managed Postgres (Neon / Supabase) + PgBouncer",
        "Hot-path indexes on rates, assignments, timesheets",
        "Read replicas for reporting load",
        "Partition timesheet_entries past ~50M rows",
    ]),
    ("Storage", [
        "Xlsx → Vercel Blob / S3",
        "DB stores only InvoiceRun.fileUrl",
        "Bucket lifecycle rule for retention",
        "No file bytes in Postgres",
    ]),
    ("Schema / Domain", [
        "New SLA / sub-cat = INSERT, zero migration",
        "Sparse rate matrix via nullable cells",
        "Multi-currency: per-account override",
        "Multi-tenant ready: scope by OrgId claim",
    ]),
]
qleft = Inches(0.5)
qtop = Inches(1.7)
qw = Inches(6.1)
qh = Inches(2.6)
qgap = Inches(0.15)
for i, (head, items) in enumerate(quads):
    r, c = divmod(i, 2)
    x = qleft + c * (qw + qgap)
    y = qtop + r * (qh + qgap)
    add_rect(s3, x, y, qw, qh, LIGHT)
    add_rect(s3, x, y, qw, Inches(0.5), NAVY)
    add_text(s3, x + Inches(0.2), y + Inches(0.08), qw - Inches(0.3), Inches(0.4),
             head, size=15, bold=True, color=WHITE)
    body = "\n".join(f"• {it}" for it in items)
    add_text(s3, x + Inches(0.2), y + Inches(0.6), qw - Inches(0.3), qh - Inches(0.7),
             body, size=12, color=DARK)

out = r"C:\Users\KaransPC\Downloads\Invoicing Project\docs\architecture-and-stack.pptx"
prs.save(out)
print(f"Saved: {out}")
